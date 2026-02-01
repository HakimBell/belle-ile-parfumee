from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
BG_MAIN = RGBColor(0x1A, 0x1D, 0x23)
BG_CARD = RGBColor(0x24, 0x28, 0x30)
BG_ACCENT = RGBColor(0x2D, 0x32, 0x3B)
BG_CODE = RGBColor(0x1E, 0x1E, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300 = RGBColor(0xCE, 0xD4, 0xDA)
GRAY_400 = RGBColor(0xAD, 0xB5, 0xBD)
GRAY_500 = RGBColor(0x8B, 0x95, 0x9E)
CORAL = RGBColor(0xFF, 0x6B, 0x6B)
PEACH = RGBColor(0xFF, 0x92, 0x7D)
GOLD = RGBColor(0xFF, 0xD9, 0x3D)
MINT = RGBColor(0x6B, 0xD4, 0xA8)
SKY = RGBColor(0x74, 0xC0, 0xFC)
LAVENDER = RGBColor(0xB1, 0x97, 0xFC)
ROSE = RGBColor(0xF7, 0x83, 0xAC)

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_MAIN

def text_box(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Arial"
    return box

def code_box(slide, left, top, width, height, code, size=10):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = code
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
    run.font.name = "Consolas"
    return box

def rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=1, radius=0.05):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius
    return shape

def glow_rect(slide, left, top, width, height, color):
    glow = rounded_rect(slide, left - 0.05, top - 0.05, width + 0.1, height + 0.1, color, None, 0, 0.08)
    glow.fill.fore_color.brightness = 0.7
    return glow

def accent_line(slide, left, top, width, color):
    line = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.adjustments[0] = 0.5
    return line

def card(slide, left, top, width, height, title, items, accent=CORAL, icon=None):
    glow_rect(slide, left, top, width, height, accent)
    rounded_rect(slide, left, top, width, height, BG_CARD, None, 0, 0.04)
    accent_line(slide, left, top, width, accent)
    title_text = f"{icon}  {title}" if icon else title
    text_box(slide, left + 0.3, top + 0.3, width - 0.6, 0.5, title_text, 16, WHITE, True)
    y = top + 0.85
    for item in items:
        if item:
            text_box(slide, left + 0.3, y, width - 0.6, 0.38, f"•  {item}", 14, GRAY_300)
        y += 0.42

def mini_card(slide, left, top, width, height, title, desc, accent=CORAL):
    glow_rect(slide, left, top, width, height, accent)
    rounded_rect(slide, left, top, width, height, BG_CARD, None, 0, 0.06)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Pt(6), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.adjustments[0] = 0.5
    text_box(slide, left + 0.25, top + 0.2, width - 0.4, 0.45, title, 15, WHITE, True)
    text_box(slide, left + 0.25, top + 0.6, width - 0.4, height - 0.7, desc, 13, GRAY_400)

def persona_card(slide, left, top, width, height, name, role, age, goals, frustrations, accent=CORAL, emoji="👤"):
    glow_rect(slide, left, top, width, height, accent)
    rounded_rect(slide, left, top, width, height, BG_CARD, None, 0, 0.04)
    accent_line(slide, left, top, width, accent)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.3), Inches(top + 0.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    text_box(slide, left + 0.3, top + 0.55, 0.8, 0.5, emoji, 28, WHITE, False, PP_ALIGN.CENTER)
    text_box(slide, left + 1.3, top + 0.4, width - 1.6, 0.4, name, 18, WHITE, True)
    text_box(slide, left + 1.3, top + 0.8, width - 1.6, 0.35, f"{role}  •  {age} ans", 12, accent)
    text_box(slide, left + 0.3, top + 1.4, width - 0.6, 0.35, "Objectifs", 13, MINT, True)
    y = top + 1.75
    for goal in goals:
        text_box(slide, left + 0.3, y, width - 0.6, 0.32, f"✓  {goal}", 12, GRAY_300)
        y += 0.35
    text_box(slide, left + 0.3, y + 0.15, width - 0.6, 0.35, "Frustrations", 13, CORAL, True)
    y += 0.5
    for frust in frustrations:
        text_box(slide, left + 0.3, y, width - 0.6, 0.32, f"✗  {frust}", 12, GRAY_400)
        y += 0.35

def slide_title(slide, title, subtitle=None, color=CORAL):
    text_box(slide, 0.6, 0.4, 12, 0.7, title, 36, color, True, PP_ALIGN.LEFT)
    accent_line(slide, 0.6, 1.1, 2.5, color)
    if subtitle:
        text_box(slide, 0.6, 1.25, 12, 0.4, subtitle, 14, GRAY_500, False, PP_ALIGN.LEFT)

def section_slide(slide, section_num, title, subtitle=None, color=CORAL):
    set_bg(slide)
    for x, y, s, c in [(-1, -0.5, 3.5, color), (11.5, 5, 3, LAVENDER), (12, -0.5, 2.5, MINT)]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
        circle.fill.solid()
        circle.fill.fore_color.rgb = c
        circle.fill.fore_color.brightness = 0.75
        circle.line.fill.background()
    rounded_rect(slide, 5.916, 2, 1.5, 0.7, color, None, 0, 0.3)
    text_box(slide, 5.916, 2.12, 1.5, 0.5, section_num, 24, WHITE, True, PP_ALIGN.CENTER)
    text_box(slide, 0, 3, 13.333, 1, title, 48, WHITE, True, PP_ALIGN.CENTER)
    accent_line(slide, 5.416, 4, 2.5, color)
    if subtitle:
        text_box(slide, 0, 4.25, 13.333, 0.5, subtitle, 18, GRAY_400, False, PP_ALIGN.CENTER)

def diagram_slide(slide, title, subtitle, image_path, color=GRAY_500):
    set_bg(slide)
    slide_title(slide, title, subtitle, color)
    rounded_rect(slide, 1.5, 1.65, 10.333, 4.7, WHITE, None, 0, 0.02)
    try:
        slide.shapes.add_picture(image_path, Inches(1.7), Inches(1.8), width=Inches(9.933))
    except:
        text_box(slide, 4, 3.5, 5, 0.5, f"[{image_path.split('/')[-1]}]", 16, GRAY_500, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITRE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

for x, y, s, c in [(-1.5, -1, 4, CORAL), (11, 5, 3.5, LAVENDER), (12, -0.5, 2, MINT)]:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
    circle.fill.solid()
    circle.fill.fore_color.rgb = c
    circle.fill.fore_color.brightness = 0.75
    circle.line.fill.background()

text_box(slide, 0, 2.2, 13.333, 1.1, "Belle Île Parfumée", 64, WHITE, True, PP_ALIGN.CENTER)
accent_line(slide, 5.166, 3.35, 3, CORAL)
text_box(slide, 0, 3.6, 13.333, 0.55, "Application E-commerce de vente de parfums", 22, GRAY_400, False, PP_ALIGN.CENTER)
rounded_rect(slide, 5.166, 4.7, 3, 0.5, CORAL, None, 0, 0.3)
text_box(slide, 5.166, 4.78, 3, 0.4, "Projet CDA 2024", 14, WHITE, True, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: SOMMAIRE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Sommaire", None, LAVENDER)

items = [
    ("01", "Présentation du projet", "Besoins, objectifs, personas, graphisme", CORAL),
    ("02", "Technologies", "Stack technique Frontend & Backend", SKY),
    ("03", "Architecture", "Architecture 3-tiers et ses avantages", MINT),
    ("04", "Conception & Modélisation", "MERISE, MCD, MLD, MPD, UML", GOLD),
    ("05", "Sécurité", "Authentification JWT, autorisation RBAC", ROSE),
    ("06", "Fonctionnalités", "Panier, code source", PEACH),
    ("07", "Démonstration", "Application en live", LAVENDER),
]

for i, (num, title, desc, color) in enumerate(items):
    y = 1.6 + i * 0.72
    rounded_rect(slide, 2.5, y, 0.55, 0.48, color, None, 0, 0.2)
    text_box(slide, 2.5, y + 0.08, 0.55, 0.35, num, 15, BG_MAIN, True, PP_ALIGN.CENTER)
    text_box(slide, 3.25, y + 0.08, 2.8, 0.4, title, 17, WHITE, True)
    text_box(slide, 6.3, y + 0.1, 5.2, 0.35, desc, 14, GRAY_400)

# ════════════════════════════════════════════════════════════════
# SECTION 1: PRÉSENTATION DU PROJET
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "01", "Présentation du Projet", "Contexte, besoins et objectifs", CORAL)

# SLIDE: BESOINS
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Besoins", "Contexte et problématique", CORAL)

card(slide, 0.6, 1.55, 5.9, 2.3, "Contexte", [
    "Auto-entrepreneur depuis 10 ans",
    "Vente de parfums grandes marques",
    "Gestion manuelle (Excel, emails)"
], SKY, "📋")

card(slide, 6.833, 1.55, 5.9, 2.3, "Problématique", [
    "Processus chronophage",
    "Risque d'erreurs humaines",
    "Mauvaise expérience client"
], CORAL, "⚠️")

card(slide, 0.6, 4.05, 12.133, 2.2, "Besoin identifié", [
    "Digitaliser l'activité avec une plateforme e-commerce moderne",
    "Automatiser la gestion des stocks et la communication client"
], MINT, "🎯")

# SLIDE: OBJECTIFS
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Objectifs", "Ce que l'application doit accomplir", GOLD)

for i, (icon, title, desc, color) in enumerate([
    ("🎯", "Automatiser", "Gestion des stocks et communication client", SKY),
    ("📦", "Centraliser", "Commandes et données en un seul endroit", LAVENDER),
    ("✨", "Améliorer l'UX", "Expérience d'achat fluide et moderne", MINT),
    ("🔒", "Sécuriser", "Protection des données avec JWT et BCrypt", ROSE),
]):
    mini_card(slide, 0.6, 1.55 + i * 1.35, 12.133, 1.1, f"{icon}  {title}", desc, color)

# SLIDE: PERSONAS
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Personas", "Nos utilisateurs cibles", ROSE)

persona_card(slide, 0.6, 1.55, 5.9, 4.8, "Sophie Martin", "Cliente régulière", 34,
    ["Acheter ses parfums préférés", "Voir les nouveautés rapidement", "Commander depuis son mobile"],
    ["Attendre les emails du dimanche", "Stock épuisé avant sa commande"], ROSE, "👩")

persona_card(slide, 6.833, 1.55, 5.9, 4.8, "Marc Dupont", "Administrateur / Vendeur", 45,
    ["Gérer son stock facilement", "Gagner du temps", "Satisfaire ses clients"],
    ["Emails manuels répétitifs", "Mises à jour Excel fastidieuses"], SKY, "👨")

# SLIDE: CHARTE GRAPHIQUE
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Charte Graphique", "Identité visuelle", ROSE)

text_box(slide, 0.6, 1.5, 4, 0.4, "PALETTE DE COULEURS", 12, ROSE, True)
for i, (name, hex_code, rgb) in enumerate([
    ("Noir", "#0A0A0A", RGBColor(0x0A, 0x0A, 0x0A)),
    ("Blanc", "#FFFFFF", RGBColor(0xFF, 0xFF, 0xFF)),
    ("Gris clair", "#F5F5F5", RGBColor(0xF5, 0xF5, 0xF5)),
    ("Gris", "#6B6B6B", RGBColor(0x6B, 0x6B, 0x6B)),
]):
    x = 0.6 + i * 1.5
    rounded_rect(slide, x, 1.95, 1.3, 0.85, rgb, GRAY_500, 1, 0.1)
    text_box(slide, x, 2.9, 1.3, 0.55, f"{name}\n{hex_code}", 10, GRAY_400, False, PP_ALIGN.CENTER)

card(slide, 6.833, 1.55, 5.9, 2.6, "Typographie", [
    "Police : Inter (Google Fonts)", "Titres : 32-44px Bold", "Corps : 16px Regular"
], SKY, "🔤")

card(slide, 6.833, 4.35, 5.9, 1.9, "Composants UI", [
    "Boutons : fond noir, hover gris", "Cartes : fond blanc, ombre légère"
], MINT, "🎨")

# SLIDE: ZONING (vide - à compléter)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Zoning", "Organisation des zones de l'interface", PEACH)
text_box(slide, 0, 3.5, 13.333, 0.5, "[Ajouter les images de zoning]", 16, GRAY_500, False, PP_ALIGN.CENTER)

# SLIDE: WIREFRAMES (vide - à compléter)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Wireframes", "Maquettes fil de fer", LAVENDER)
text_box(slide, 0, 3.5, 13.333, 0.5, "[Ajouter les wireframes]", 16, GRAY_500, False, PP_ALIGN.CENTER)

# SLIDE: MAQUETTES (vide - à compléter)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Maquettes", "Design final de l'interface", MINT)
text_box(slide, 0, 3.5, 13.333, 0.5, "[Ajouter les maquettes]", 16, GRAY_500, False, PP_ALIGN.CENTER)

# SLIDE: ARBORESCENCE
slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "Arborescence", "Structure de navigation du site",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/arborescence_diagram.png", SKY)

# ════════════════════════════════════════════════════════════════
# SECTION 2: TECHNOLOGIES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "02", "Technologies", "Stack technique", SKY)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Stack Technique", "Technologies utilisées", SKY)

card(slide, 0.6, 1.55, 5.9, 4.7, "Frontend", [
    "React 19 + TypeScript 5.9", "Vite 7 (build tool)", "React Router DOM (navigation)",
    "Axios (requêtes HTTP)", "Context API (état global)"
], SKY, "⚛️")

card(slide, 6.833, 1.55, 5.9, 4.7, "Backend", [
    "Spring Boot 3.5 + Java 21", "Spring Security + JWT", "Spring Data JPA (ORM)",
    "PostgreSQL (base de données)", "JavaMail (emails)"
], MINT, "☕")

# ════════════════════════════════════════════════════════════════
# SECTION 3: ARCHITECTURE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "03", "Architecture", "Architecture 3-tiers", MINT)

slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "Architecture Globale", "Architecture 3-tiers",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/architecture_diagram.png", MINT)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Avantages de l'Architecture 3-Tiers", None, MINT)

for i, (icon, title, desc, color) in enumerate([
    ("🔀", "Séparation des responsabilités", "Chaque couche a un rôle précis et distinct", SKY),
    ("📈", "Scalabilité", "Possibilité de faire évoluer chaque couche indépendamment", LAVENDER),
    ("🔧", "Maintenabilité", "Modifications isolées sans impact sur les autres couches", MINT),
    ("🧪", "Testabilité", "Tests unitaires et d'intégration facilités", GOLD),
]):
    mini_card(slide, 0.6, 1.55 + i * 1.35, 12.133, 1.1, f"{icon}  {title}", desc, color)

# ════════════════════════════════════════════════════════════════
# SECTION 4: CONCEPTION & MODÉLISATION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "04", "Conception & Modélisation", "MERISE et UML", GOLD)

# MERISE
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Méthode MERISE", "Modélisation de la base de données", GOLD)

card(slide, 0.6, 1.55, 12.133, 4.7, "Les étapes de la méthode MERISE", [
    "MCD (Modèle Conceptuel de Données) : Représentation des entités et associations",
    "MLD (Modèle Logique de Données) : Traduction en tables relationnelles",
    "MPD (Modèle Physique de Données) : Implémentation SQL spécifique au SGBD",
    "", "Cette méthode permet une conception rigoureuse et progressive de la base de données"
], GOLD, "📊")

# MCD
slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "MCD", "Modèle Conceptuel de Données",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/mcd_diagram.png", GOLD)

# MLD
slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "MLD", "Modèle Logique de Données",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/mld_diagram.png", PEACH)

# MPD
slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "MPD", "Modèle Physique de Données - Script SQL",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/mpd_diagram.png", CORAL)

# UML Intro
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Diagrammes UML", "Unified Modeling Language", LAVENDER)

card(slide, 0.6, 1.55, 3.8, 4.7, "Diagrammes utilisés", [
    "Cas d'utilisation", "Diagramme de classes", "Diagramme de séquence"
], LAVENDER, "📐")

card(slide, 4.6, 1.55, 8.133, 4.7, "Pourquoi UML ?", [
    "Langage de modélisation standardisé et universel",
    "Facilite la communication entre développeurs",
    "Documente l'architecture du système",
    "Visualise la structure avant l'implémentation",
    "Indispensable pour la conception orientée objet"
], SKY, "💡")

# Use Case
slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "Diagramme de Cas d'Utilisation", "Interactions utilisateur-système",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/usecase_diagram.png", LAVENDER)

# Classes
slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "Diagramme de Classes", "Relations : composition et agrégation",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/classes_diagram.png", GOLD)

# Séquence
slide = prs.slides.add_slide(prs.slide_layouts[6])
diagram_slide(slide, "Diagramme de Séquence", "Ajout d'un produit au panier",
    "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/diagrams/sequence_diagram.png", MINT)

# ════════════════════════════════════════════════════════════════
# SECTION 5: SÉCURITÉ
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "05", "Sécurité", "Authentification et autorisation", ROSE)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Sécurité", "Protection de l'application", ROSE)

card(slide, 0.6, 1.55, 5.9, 4.7, "Authentification", [
    "JWT (JSON Web Token)", "Token stocké en localStorage", "Expiration : 24 heures",
    "BCrypt (hashage mots de passe)", "Validation des entrées"
], SKY, "🔐")

card(slide, 6.833, 1.55, 5.9, 4.7, "Autorisation", [
    "Rôles : CLIENT, ADMIN", "Routes protégées (PrivateRoute)", "Filtre JWT côté backend",
    "CORS configuré", "@PreAuthorize sur endpoints"
], ROSE, "🛡️")

# ════════════════════════════════════════════════════════════════
# SECTION 6: FONCTIONNALITÉS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "06", "Fonctionnalités", "Features clés de l'application", PEACH)

# Panier
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Gestion du Panier", "Double persistance selon connexion", MINT)

card(slide, 0.6, 1.55, 5.9, 2.1, "Mode Non Connecté", [
    "Stockage localStorage", "Persistance côté client"
], GRAY_500, "💾")

card(slide, 6.833, 1.55, 5.9, 2.1, "Mode Connecté", [
    "Stockage en BDD (Order PENDING)", "Synchronisation temps réel"
], SKY, "☁️")

card(slide, 0.6, 3.85, 12.133, 2.4, "Caractéristiques clés", [
    "Fusion automatique du panier local à la connexion",
    "Vérification du stock en temps réel avant validation",
    "Email de confirmation automatique après commande"
], MINT, "✨")

# CODE SLIDE
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Code : Ajout au Panier", "CartContext.tsx - Fonction addToCart", SKY)

# Code background
rounded_rect(slide, 0.5, 1.55, 12.333, 4.8, BG_CODE, GRAY_500, 1, 0.02)

# File indicator
rounded_rect(slide, 0.5, 1.55, 3, 0.4, BG_ACCENT, None, 0, 0.15)
text_box(slide, 0.7, 1.6, 2.6, 0.3, "📄 CartContext.tsx", 11, GRAY_300)

code = '''const addToCart = useCallback(async (product: Product, quantity: number) => {
  setError(null);
  const isAuth = getIsAuthenticated();
  const email = getUserEmail();

  if (isAuth && email) {
    // Mode connecté : API Backend
    setIsLoading(true);
    try {
      const updatedCart = await cartApiService.addItem(email, {
        productCode: product.productCode,
        quantity
      });
      setItems(convertBackendItems(updatedCart.items));
    } catch (err) {
      setError('Erreur lors de l\\'ajout au panier');
    } finally {
      setIsLoading(false);
    }
  } else {
    // Mode non connecté : localStorage
    setItems(prevItems => {
      const existingItem = prevItems.find(
        item => item.product.productCode === product.productCode
      );
      if (existingItem) {
        return prevItems.map(item =>
          item.product.productCode === product.productCode
            ? { ...item, quantity: Math.min(item.quantity + quantity, product.stock) }
            : item
        );
      }
      return [...prevItems, { product, quantity }];
    });
  }
}, []);'''

code_box(slide, 0.7, 2.05, 11.933, 4.2, code, 9)

# Annotations
rounded_rect(slide, 9.5, 2.3, 3.2, 0.5, SKY, None, 0, 0.2)
text_box(slide, 9.6, 2.35, 3, 0.4, "→ Vérification auth", 10, WHITE, True)

rounded_rect(slide, 9.5, 3.8, 3.2, 0.5, MINT, None, 0, 0.2)
text_box(slide, 9.6, 3.85, 3, 0.4, "→ Appel API Backend", 10, WHITE, True)

rounded_rect(slide, 9.5, 5.3, 3.2, 0.5, PEACH, None, 0, 0.2)
text_box(slide, 9.6, 5.35, 3, 0.4, "→ Fallback localStorage", 10, WHITE, True)

# Code produit
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Génération du Code Produit", "Identifiant unique automatique", LAVENDER)

rounded_rect(slide, 3, 1.55, 7.333, 0.6, BG_CARD, LAVENDER, 2, 0.15)
text_box(slide, 3, 1.62, 7.333, 0.5, "XXX – TTT – MMM – YYYY", 22, LAVENDER, True, PP_ALIGN.CENTER)

for i, (code, label, example, color) in enumerate([
    ("XXX", "Marque", "CHA, DIO", SKY),
    ("TTT", "Concentration", "EDP, EDT", LAVENDER),
    ("MMM", "Taille (ml)", "050, 100", MINT),
    ("YYYY", "Séquentiel", "0001", GOLD),
]):
    x = 0.8 + i * 3
    glow_rect(slide, x, 2.4, 2.7, 1.4, color)
    rounded_rect(slide, x, 2.4, 2.7, 1.4, BG_CARD, None, 0, 0.06)
    accent_line(slide, x, 2.4, 2.7, color)
    text_box(slide, x, 2.55, 2.7, 0.45, code, 20, color, True, PP_ALIGN.CENTER)
    text_box(slide, x, 3, 2.7, 0.35, label, 13, WHITE, False, PP_ALIGN.CENTER)
    text_box(slide, x, 3.35, 2.7, 0.35, example, 12, GRAY_400, False, PP_ALIGN.CENTER)

rounded_rect(slide, 2.5, 4.1, 8.333, 1.3, BG_ACCENT, None, 0, 0.06)
text_box(slide, 2.5, 4.35, 8.333, 0.45, "Exemple :  CHA-EDP-100-0001", 18, WHITE, True, PP_ALIGN.CENTER)
text_box(slide, 2.5, 4.8, 8.333, 0.4, "Chanel N°5 Eau de Parfum 100ml", 14, GRAY_400, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SECTION 7: DÉMONSTRATION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_slide(slide, "07", "Démonstration", "Application en live", LAVENDER)

# CONCLUSION
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
slide_title(slide, "Conclusion", "Bilan et perspectives", MINT)

card(slide, 0.6, 1.55, 5.9, 4.7, "Réalisations", [
    "Frontend React 19 moderne", "Backend Spring Boot 3.5 robuste", "Base de données PostgreSQL",
    "Sécurité JWT + BCrypt", "Code maintenable et documenté"
], MINT, "✅")

card(slide, 6.833, 1.55, 5.9, 4.7, "Évolutions possibles", [
    "Paiement en ligne (Stripe)", "Tests E2E (Cypress)", "CI/CD (GitHub Actions)",
    "PWA (Progressive Web App)", "Application mobile native"
], SKY, "🚀")

# QUESTIONS
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

for x, y, s, c in [(-0.5, 5, 3, ROSE), (11.5, -0.5, 3, SKY), (-1, -1, 2.5, GOLD)]:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
    circle.fill.solid()
    circle.fill.fore_color.rgb = c
    circle.fill.fore_color.brightness = 0.75
    circle.line.fill.background()

text_box(slide, 0, 2.4, 13.333, 1, "Questions ?", 54, CORAL, True, PP_ALIGN.CENTER)
accent_line(slide, 5.416, 3.5, 2.5, CORAL)
text_box(slide, 0, 3.75, 13.333, 0.5, "Merci de votre attention", 18, GRAY_400, False, PP_ALIGN.CENTER)

# SAVE
output = "/Users/hakim/Desktop/CDA/belle-ile-parfumee/docs/presentation-cda.pptx"
prs.save(output)
print("✅ Présentation créée avec succès!")
print(f"📍 {output}")
print("📊 31 slides - Avec slides vides pour Zoning, Wireframes et Maquettes")